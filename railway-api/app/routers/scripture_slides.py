"""
Scripture slides router for generating Bible verse PowerPoint presentations
"""
import os
from typing import List, Dict
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from app.core.schemas import ScriptureSlideRequest
from app.core.files import create_temp_file
from .slides.utils import (
    # Base presentation functions
    create_presentation,
    save_presentation,
    process_background_image,
    set_slide_background,
    cleanup_temp_file,
    # Text effect functions
    add_text_glow
)

router = APIRouter()


def _get_book_name(book_code: str, is_tongan: bool = False) -> str:
    """Convert book code to readable book name
    
    Args:
        book_code: The book code (e.g., 'JHN', 'MAT')
        is_tongan: If True, return Tongan book name for TMB translation
    """
    if is_tongan:
        # Tongan book names
        tongan_book_names = {
            'GEN': 'Kenesi', 'EXO': 'ʻEkisotosi', 'LEV': 'Levitiko', 'NUM': 'Nemipia', 'DEU': 'Teutelonomi',
            'JOS': 'Siosiua', 'JDG': 'Fakamaau', 'RUT': 'Lute', '1SA': '1 Samueli', '2SA': '2 Samueli',
            '1KI': '1 Tuʻi', '2KI': '2 Tuʻi', '1CH': '1 Kalonikali', '2CH': '2 Kalonikali',
            'EZR': 'ʻEsila', 'NEH': 'Nehimaia', 'EST': 'ʻEseta', 'JOB': 'Siope', 'PSA': 'Saame',
            'PRO': 'Lea Fakatatauki', 'ECC': 'ʻEkilisiasitesi', 'SNG': 'Hiva ʻa Solomone', 'ISA': 'ʻAisea',
            'JER': 'Selemaia', 'LAM': 'Tangilāulau', 'EZK': 'ʻIsikieli', 'DAN': 'Taniela',
            'HOS': 'Hosea', 'JOL': 'Soeli', 'AMO': 'ʻAmosi', 'OBA': 'ʻOpataia', 'JON': 'Siona',
            'MIC': 'Maika', 'NAM': 'Nahumi', 'HAB': 'Hapakuki', 'ZEP': 'Sefanaia', 'HAG': 'Hakai',
            'ZEC': 'Sekalaia', 'MAL': 'Malaki',
            'MAT': 'Matiu', 'MRK': 'Maʻake', 'LUK': 'Luke', 'JHN': 'Sione', 'ACT': 'Ngaue',
            'ROM': 'Loma', '1CO': '1 Kolinito', '2CO': '2 Kolinito', 'GAL': 'Kalātia',
            'EPH': 'ʻEfeso', 'PHP': 'Filipai', 'COL': 'Kolosi', '1TH': '1 Tesalonaika',
            '2TH': '2 Tesalonaika', '1TI': '1 Timote', '2TI': '2 Timote', 'TIT': 'Taitosi',
            'PHM': 'Filimona', 'HEB': 'Hepelū', 'JAS': 'Semisi', '1PE': '1 Pita', '2PE': '2 Pita',
            '1JN': '1 Sione', '2JN': '2 Sione', '3JN': '3 Sione', 'JUD': 'Siuta', 'REV': 'Fakahā'
        }
        return tongan_book_names.get(book_code, book_code)
    else:
        # English book names
        book_names = {
            'GEN': 'Genesis', 'EXO': 'Exodus', 'LEV': 'Leviticus', 'NUM': 'Numbers', 'DEU': 'Deuteronomy',
            'JOS': 'Joshua', 'JDG': 'Judges', 'RUT': 'Ruth', '1SA': '1 Samuel', '2SA': '2 Samuel',
            '1KI': '1 Kings', '2KI': '2 Kings', '1CH': '1 Chronicles', '2CH': '2 Chronicles',
            'EZR': 'Ezra', 'NEH': 'Nehemiah', 'EST': 'Esther', 'JOB': 'Job', 'PSA': 'Psalms',
            'PRO': 'Proverbs', 'ECC': 'Ecclesiastes', 'SNG': 'Song of Songs', 'ISA': 'Isaiah',
            'JER': 'Jeremiah', 'LAM': 'Lamentations', 'EZK': 'Ezekiel', 'DAN': 'Daniel',
            'HOS': 'Hosea', 'JOL': 'Joel', 'AMO': 'Amos', 'OBA': 'Obadiah', 'JON': 'Jonah',
            'MIC': 'Micah', 'NAM': 'Nahum', 'HAB': 'Habakkuk', 'ZEP': 'Zephaniah', 'HAG': 'Haggai',
            'ZEC': 'Zechariah', 'MAL': 'Malachi',
            'MAT': 'Matthew', 'MRK': 'Mark', 'LUK': 'Luke', 'JHN': 'John', 'ACT': 'Acts',
            'ROM': 'Romans', '1CO': '1 Corinthians', '2CO': '2 Corinthians', 'GAL': 'Galatians',
            'EPH': 'Ephesians', 'PHP': 'Philippians', 'COL': 'Colossians', '1TH': '1 Thessalonians',
            '2TH': '2 Thessalonians', '1TI': '1 Timothy', '2TI': '2 Timothy', 'TIT': 'Titus',
            'PHM': 'Philemon', 'HEB': 'Hebrews', 'JAS': 'James', '1PE': '1 Peter', '2PE': '2 Peter',
            '1JN': '1 John', '2JN': '2 John', '3JN': '3 John', 'JUD': 'Jude', 'REV': 'Revelation'
        }
        return book_names.get(book_code, book_code)


def create_scripture_slides(
    reference: Dict[str, str],
    verses: List[Dict[str, str]],
    output_file: str = "scripture_slides.pptx",
    background_image: str | None = None,
    verses_alt: List[Dict[str, str]] | None = None,
) -> str:
    """Create scripture slides.

    reference: { 'book': 'PS', 'chapter': 23 }
    verses: list of { 'verse': int, 'text': str } - primary translation (NRSVUE)
    verses_alt: optional list of { 'verse': int, 'text': str } - alternate translation (TMB) for combined mode
    """
    prs = create_presentation()

    # Process background image - always expect base64
    background_image_path = None
    if background_image:
        background_image_path = process_background_image(background_image)

    # No default background - frontend should always provide one

    slide_count = 0

    # If verses_alt is provided, we're in combined mode - alternate between translations
    if verses_alt:
        # Create a mapping of verse numbers to both translations
        verse_map = {}
        for v in verses:
            verse_num = v.get("verse")
            if verse_num:
                verse_map[verse_num] = {"nrsvue": v.get("text", "").strip(), "tmb": ""}
        
        for v in verses_alt:
            verse_num = v.get("verse")
            if verse_num and verse_num in verse_map:
                verse_map[verse_num]["tmb"] = v.get("text", "").strip()
        
        # Generate slides alternating between translations for each verse
        for verse_num in sorted(verse_map.keys()):
            texts = verse_map[verse_num]
            
            # First slide: NRSVUE
            if texts["nrsvue"]:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                set_slide_background(slide, background_image_path)
                
                # Add scripture placeholder image in top right corner
                placeholder_path = os.path.join(
                    os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
                    'public', 'images', 'placeholders', 'scripture-placeholder.jpg'
                )
                
                if os.path.exists(placeholder_path):
                    # Position: 9.33 inches horizontal, 0.25 inches vertical
                    left = Inches(9.33)
                    top = Inches(0.25)
                    # Size: width 3.72 inches, height 2 inches
                    width = Inches(3.72)
                    height = Inches(2)
                    slide.shapes.add_picture(placeholder_path, left, top, width=width, height=height)
                
                _add_verse_content(slide, reference['book'], reference['chapter'], verse_num, texts["nrsvue"], "NRSVUE")
                slide_count += 1
            
            # Second slide: TMB
            if texts["tmb"]:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                set_slide_background(slide, background_image_path)
                
                # Add scripture placeholder image in top right corner
                placeholder_path = os.path.join(
                    os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
                    'public', 'images', 'placeholders', 'scripture-placeholder.jpg'
                )
                
                if os.path.exists(placeholder_path):
                    # Position: 9.33 inches horizontal, 0.25 inches vertical
                    left = Inches(9.33)
                    top = Inches(0.25)
                    # Size: width 3.72 inches, height 2 inches
                    width = Inches(3.72)
                    height = Inches(2)
                    slide.shapes.add_picture(placeholder_path, left, top, width=width, height=height)
                
                _add_verse_content(slide, reference['book'], reference['chapter'], verse_num, texts["tmb"], "TMB")
                slide_count += 1
    else:
        # Single translation mode - use verses as before
        for v in verses:
            verse_num = v.get("verse")
            text = v.get("text", "").strip()
            if not text:
                continue

            # Create exactly one slide per verse
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            set_slide_background(slide, background_image_path)
            
            # Add scripture placeholder image in top right corner
            placeholder_path = os.path.join(
                os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
                'public', 'images', 'placeholders', 'scripture-placeholder.jpg'
            )
            
            if os.path.exists(placeholder_path):
                # Position: 9.33 inches horizontal, 0.25 inches vertical
                left = Inches(9.33)
                top = Inches(0.25)
                # Size: width 3.72 inches, height 2 inches
                width = Inches(3.72)
                height = Inches(2)
                slide.shapes.add_picture(placeholder_path, left, top, width=width, height=height)
            
            _add_verse_content(slide, reference['book'], reference['chapter'], verse_num, text, None)
            slide_count += 1

    if background_image:
        cleanup_temp_file(background_image_path)

    save_presentation(prs, output_file)
    return f"Created {slide_count} scripture slides"


def _add_verse_content(slide, book, chapter, verse_num, text, translation_label=None):
    """Helper function to add verse content to a slide"""
    # Title: book, chapter, verse number with optional translation label
    title_left = Inches(0.26)
    title_top = Inches(0.21)
    title_width = Inches(8.91)
    title_height = Inches(1.2)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)

    title_para = title_frame.paragraphs[0]
    # Include book, chapter, and verse
    # Use Tongan book name for TMB translation
    is_tongan = translation_label == "TMB"
    book_name = _get_book_name(book, is_tongan)
    title_para.text = f"{book_name} {chapter}:{verse_num}"
    title_para.font.name = "Arial Narrow"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.underline = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)
    title_para.alignment = PP_ALIGN.LEFT
    add_text_glow(title_para, glow_radius=6, color_rgb=(255, 255, 255))
    
    # Add translation label on a new line if provided (combined mode only)
    if translation_label:
        translation_para = title_frame.add_paragraph()
        translation_para.text = translation_label
        translation_para.font.name = "Arial Narrow"
        translation_para.font.size = Pt(36)
        translation_para.font.bold = True
        translation_para.font.color.rgb = RGBColor(0, 0, 0)
        translation_para.alignment = PP_ALIGN.LEFT
        add_text_glow(translation_para, glow_radius=6, color_rgb=(255, 255, 255))

    # Content - display entire verse text on one slide
    content_top = Inches(1.7)
    content_left = Inches(0)
    content_width = Inches(13.33)
    content_height = Inches(5.0)

    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
    content_frame.margin_left = Inches(0.5)
    content_frame.margin_right = Inches(0.5)
    content_frame.clear()

    # Add the entire verse text as a single paragraph
    p = content_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial Narrow"
    # Dynamically adjust font size based on text length
    if len(text) > 300:
        p.font.size = Pt(36)
    elif len(text) > 200:
        p.font.size = Pt(44)
    elif len(text) > 100:
        p.font.size = Pt(52)
    else:
        p.font.size = Pt(58)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    add_text_glow(p, glow_radius=6, color_rgb=(255, 255, 255))


@router.post("/generate-scripture-slides")
async def generate_scripture_slides_endpoint(request: ScriptureSlideRequest):
    """Generate scripture slides"""
    try:
        # Create temporary file for the PowerPoint
        output_path = create_temp_file(suffix='.pptx')
        
        # Use background image if provided (base64 only)
        background = request.background_image
        
        # Generate the PowerPoint
        create_scripture_slides(
            reference=request.reference,
            verses=request.verses,
            output_file=output_path,
            background_image=request.background_image,
            verses_alt=request.verses_alt
        )
        
        # Return the file
        return FileResponse(
            path=output_path,
            filename="scripture.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))