"""
Hymn slides router for generating hymn PowerPoint presentations
"""
import os
import re
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any
from app.core.schemas import HymnRequest
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from app.core.files import create_temp_file
from .slides.utils import (
    # Base presentation functions
    create_presentation, 
    save_presentation, 
    process_background_image,
    set_slide_background,
    cleanup_temp_file,
    # Text effect functions
    add_text_glow
)

router = APIRouter()


def format_hymn_number(hymnal, hymn_number):
    """Format hymn number with hymnal abbreviation."""
    hymnal_names = {
        'umh': 'UMH',
        'fws': 'FWS',
        'thb': 'THB'
    }
    return f"{hymnal_names.get(hymnal.lower(), hymnal.upper())} {hymn_number}"


def format_copyright_info(hymn_data):
    """Format copyright information, handling empty fields gracefully."""
    copyright_parts = []
    
    # Add text copyright if available
    if hymn_data.get('text_copyright') and hymn_data['text_copyright'].strip():
        copyright_parts.append(f"Text: {hymn_data['text_copyright']}")
    
    # Add tune copyright if available
    if hymn_data.get('tune_copyright') and hymn_data['tune_copyright'].strip():
        copyright_parts.append(f"Music: {hymn_data['tune_copyright']}")
    
    # If no copyright info, check if it might be public domain
    if not copyright_parts:
        # Look for dates in author/composer info to guess if it's likely public domain
        author = hymn_data.get('author', '')
        composer = hymn_data.get('composer', '')
        
        # Simple heuristic: if we have dates and they're old, assume public domain
        if author or composer:
            # Check for old dates that would indicate public domain
            all_text = str(author) + " " + str(composer)
            if any(str(year) in all_text for year in range(1800, 1923)):
                copyright_parts.append("Public Domain")
    
    return copyright_parts


def add_hymn_cover_slide(prs, hymn_data, background_image_path=None):
    """
    Add a hymn cover slide to an existing presentation.
    
    Args:
        prs: PowerPoint presentation object
        hymn_data: Dict with hymn data containing title, hymn_number, hymnal, etc.
        background_image_path: Path to background image (optional)
    
    Returns:
        The slide object that was created
    """
    # Get the blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background image
    set_slide_background(slide, background_image_path)
    
    
    # Add main title - large, centered, with glow effect
    title_left = Inches(0)
    title_top = Inches(1.88)  # 1.88" from top
    title_width = prs.slide_width  # Full width
    title_height = Inches(3.33)  # 3.33" high
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    
    title_run = title_p.runs[0] if title_p.runs else title_p.add_run()
    title_run.text = f'"{hymn_data["title"]}"'
    
    # Style the title to match reference
    title_font = title_run.font
    title_font.name = 'Arial Narrow'
    title_font.size = Pt(66)  # Large size for main title
    title_font.bold = True
    title_font.color.rgb = RGBColor(0, 0, 0)  # Black
    
    # Add white glow effect to title
    add_text_glow(title_p, glow_radius=6, color_rgb=(255, 255, 255))
    
    # Add hymn information box - smaller text, centered
    info_left = Inches(0)
    info_top = Inches(5.17)  # 5.17" from top
    info_width = prs.slide_width  # Full width
    info_height = Inches(2.33)  # 2.33" high
    
    info_box = slide.shapes.add_textbox(info_left, info_top, info_width, info_height)
    info_frame = info_box.text_frame
    info_frame.clear()
    info_frame.word_wrap = True
    info_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically
    info_frame.margin_left = Inches(0.5)
    info_frame.margin_right = Inches(0.5)
    
    # Build the information lines
    info_lines = []
    
    # Hymn number
    hymn_num = format_hymn_number(hymn_data['hymnal'], hymn_data['hymn_number'])
    info_lines.append(hymn_num)
    
    # Author (Words)
    if hymn_data.get('author') and hymn_data['author'].strip():
        info_lines.append(f"Words: {hymn_data['author']}")
    
    # Composer (Music)
    if hymn_data.get('composer') and hymn_data['composer'].strip():
        info_lines.append(f"Music: {hymn_data['composer']}")
    
    # Tune name
    if hymn_data.get('tune_name') and hymn_data['tune_name'].strip():
        info_lines.append(f"Tune: {hymn_data['tune_name']}")
    
    # Copyright information
    copyright_parts = format_copyright_info(hymn_data)
    info_lines.extend(copyright_parts)
    
    # Add each line as a separate paragraph
    for i, line in enumerate(info_lines):
        if i == 0:
            p = info_frame.paragraphs[0]
        else:
            p = info_frame.add_paragraph()
        
        p.alignment = PP_ALIGN.CENTER
        
        run = p.runs[0] if p.runs else p.add_run()
        run.text = line
        
        # Style the information text to match reference
        font = run.font
        font.name = 'Arial Narrow'
        font.size = Pt(24)  # Medium size for info
        font.bold = True
        font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # Add white glow effect to info text
        add_text_glow(p, glow_radius=4, color_rgb=(255, 255, 255))
    
    return slide


def create_hymn_slides(hymn_data, output_file='hymn_slides.pptx', background_image=None, include_cover=True):
    """
    Create hymn slides from hymn data with parsed lyrics.
    hymn_data should have keys: 'hymn_number', 'title', 'hymnal', 'lyrics'
    where lyrics is a list of objects with 'page_name' and 'text' keys.
    
    Args:
        hymn_data: Dict containing hymn information
        output_file: Path to save the presentation
        background_image: Background image (file path or base64 data)
        include_cover: Whether to include a cover slide (default: True)
    """
    # Create a presentation object
    prs = create_presentation()
    
    # Process background image - always expect base64
    background_image_path = None
    if background_image:
        background_image_path = process_background_image(background_image)
    
    slide_count = 0
    
    # Add cover slide if requested
    if include_cover:
        add_hymn_cover_slide(prs, hymn_data, background_image_path)
        slide_count += 1
    
    # Process each verse/section
    for verse_idx, verse_data in enumerate(hymn_data['lyrics']):
        page_name = verse_data.get('page_name', '')
        verse_text = verse_data.get('text', '')
        
        # Split verse text by <br> tags to create individual slides
        slides = verse_text.split('<br>')
        
        for slide_idx, slide_text in enumerate(slides):
            slide_text = slide_text.strip()
            if not slide_text:
                continue
                
            slide_count += 1
            add_hymn_slide(
                prs, 
                hymn_data, 
                slide_text, 
                page_name,
                verse_idx + 1,  # Verse number (1-based)
                slide_idx + 1,  # Slide within verse (1-based)
                len(slides),    # Total slides in this verse
                background_image_path
            )
    
    # Clean up temporary background image file if created
    if background_image:
        cleanup_temp_file(background_image_path)
    
    # Save the presentation
    save_presentation(prs, output_file)
    
    # Create descriptive message
    if include_cover:
        hymn_slide_count = slide_count - 1  # Subtract the cover slide
        return f"Created {slide_count} slides for {hymn_data['title']} (1 cover + {hymn_slide_count} hymn slides)"
    else:
        return f"Created {slide_count} hymn slides for {hymn_data['title']}"


def add_hymn_slide(prs, hymn_data, slide_text, page_name, verse_num, slide_in_verse, total_in_verse, background_image_path=None):
    """Add a hymn slide with formatting and page names"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_slide_background(slide, background_image_path)
    
    # Add hymn placeholder image in top right corner
    placeholder_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
        'public', 'images', 'placeholders', 'hymn-placeholder.jpg'
    )
    
    if os.path.exists(placeholder_path):
        # Position: 9.33 inches horizontal, 0.25 inches vertical
        left = Inches(9.33)
        top = Inches(0.25)
        # Size: width 3.72 inches, height 2 inches
        width = Inches(3.72)
        height = Inches(2)
        slide.shapes.add_picture(placeholder_path, left, top, width=width, height=height)
    
    # Add title with hymn information and verse indicator
    title_left = Inches(0.26)
    title_top = Inches(0.21)
    title_width = Inches(8.91)
    title_height = Inches(2.04)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    
    # First paragraph - title (converted to title case)
    title_para = title_frame.paragraphs[0]
    title_para.text = f"{hymn_data['title'].title()}"
    title_para.font.name = "Arial Narrow"
    title_para.font.size = Pt(50)
    title_para.font.bold = True
    title_para.font.underline = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add white glow effect to title
    add_text_glow(title_para, glow_radius=6, color_rgb=(255, 255, 255))
    
    # Add verse/page indicator in the same text box if present
    if page_name:
        # Add a new paragraph for the verse indicator
        verse_para = title_frame.add_paragraph()
        
        # Show only the page name (e.g., "Verse 1"), without slide position
        verse_para.text = page_name
        
        verse_para.font.name = "Arial Narrow"
        verse_para.font.size = Pt(40)
        verse_para.font.bold = True
        verse_para.font.color.rgb = RGBColor(0, 0, 0)
        verse_para.alignment = PP_ALIGN.LEFT
        
        # Add white glow effect to verse indicator
        add_text_glow(verse_para, glow_radius=6, color_rgb=(255, 255, 255))
    
    # Position lyrics content
    content_top = Inches(2.4)  # Below title
    content_left = Inches(0)  # Full width
    content_width = Inches(13.33)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.TOP
    content_frame.margin_left = Inches(0.5)
    content_frame.margin_right = Inches(0.5)
    
    # Clear default paragraph
    content_frame.clear()
    
    # Add lyrics text line by line
    lines = slide_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        # Set the text
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        
        # Apply formatting to the paragraph
        p.font.name = "Arial Narrow"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add glow effect to the paragraph
        add_text_glow(p, glow_radius=6, color_rgb=(255, 255, 255))


@router.post("/generate-hymn-slides")
async def generate_hymn_slides_endpoint(data: Dict[str, Any]):
    """Generate hymn PowerPoint slides"""
    try:
        # Check if hymn data is nested under 'hymn' key (from frontend)
        if 'hymn' in data:
            hymn_request = data['hymn']
            title = hymn_request.get('title', 'Hymn')
            number = hymn_request.get('number', '')
            hymnal = hymn_request.get('hymnal', 'UMH').lower()
            # For nested format, lyrics come with the hymn
            lyrics_from_request = hymn_request.get('lyrics', [])
        else:
            # Direct format (from test scripts)
            title = data.get('title', 'Hymn')
            number = data.get('number', '')
            hymnal = data.get('hymnal', 'UMH').lower()
            lyrics_from_request = data.get('lyrics', [])
        
        background_image = data.get('background_image')
        
        # Convert verses to lyrics format if needed
        def convert_verses_to_lyrics(verses):
            """Convert verses array to lyrics format"""
            if not verses:
                return []
            
            lyrics = []
            for i, verse in enumerate(verses):
                if isinstance(verse, dict):
                    # Already in structured format
                    lyrics.append({
                        'page_name': verse.get('page_name', f'Verse {i+1}'),
                        'text': verse.get('text', '')
                    })
                else:
                    # Simple string format
                    lyrics.append({
                        'page_name': f'Verse {i+1}',
                        'text': str(verse)
                    })
            return lyrics
        
        # All hymn data must come from the request
        if 'hymn' in data:
            # Nested format from frontend
            hymn_request = data['hymn']
            hymn_info = {
                'title': hymn_request.get('title', 'Hymn'),
                'hymn_number': hymn_request.get('number', ''),
                'hymnal': hymn_request.get('hymnal', 'UMH').lower(),
                'author': hymn_request.get('author', ''),
                'composer': hymn_request.get('composer', ''),
                'tune_name': hymn_request.get('tune_name', ''),
                'text_copyright': hymn_request.get('text_copyright', ''),
                'tune_copyright': hymn_request.get('tune_copyright', ''),
                'lyrics': hymn_request.get('lyrics', [])
            }
        else:
            # Direct format
            hymn_info = {
                'title': data.get('title', 'Hymn'),
                'hymn_number': data.get('number', ''),
                'hymnal': data.get('hymnal', 'UMH').lower(),
                'author': data.get('author', ''),
                'composer': data.get('composer', ''),
                'tune_name': data.get('tune_name', ''),
                'text_copyright': data.get('text_copyright', data.get('copyright', '')),
                'tune_copyright': data.get('tune_copyright', ''),
                'lyrics': data.get('lyrics', convert_verses_to_lyrics(data.get('verses', [])))
            }
        
        # Create temporary file for the PowerPoint
        output_path = create_temp_file(suffix='.pptx')
        
        # Use background image if provided (base64 only)
        background = background_image
        
        # Generate the PowerPoint
        result = create_hymn_slides(hymn_info, output_path, background)
        
        # Return the file
        return FileResponse(
            path=output_path,
            filename=f"hymn_{hymnal}_{number}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))