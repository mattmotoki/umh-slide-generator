"""
Call to Worship slides router for generating responsive reading PowerPoint presentations
"""
import os
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from typing import Dict, Any
from app.core.schemas import CallToWorshipRequest
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from app.core.files import create_temp_file
from .slides.utils import (
    # Base presentation functions
    create_presentation, 
    save_presentation, 
    process_background_image,
    set_slide_background,
    cleanup_temp_file,
    # Text effect functions
    add_text_glow,
    set_hanging_indent,
    add_text_highlight,
    add_text_outline,
    add_text_glow_scheme,
    add_end_paragraph_glow_and_highlight,
    add_run_glow,
    add_run_highlight,
    add_run_fill,
    set_run_highlight,
    add_run_outline,
    set_run_effects
)

router = APIRouter()


def calculate_font_size(text, max_size=36, min_size=26):
    """Calculate appropriate font size based on text length for combined text box"""
    # This function is no longer used, keeping for backward compatibility
    return 50  # Fixed font size


def create_call_to_worship_slides_from_dict(pairs_list, output_file='call_to_worship.pptx', background_image=None):
    """
    Create Call to Worship slides from a list of dictionaries.
    Each dictionary should have 'Leader' and 'People' keys.
    Optional background_image should be base64 encoded string.
    """
    # Create a presentation object
    prs = create_presentation()
    
    # Process background image if provided, otherwise use default
    # Process background image - always expect base64
    background_image_path = None
    if background_image:
        background_image_path = process_background_image(background_image)
    
    # Create a slide for each pair
    for pair in pairs_list:
        add_call_to_worship_slide(prs, pair['Leader'], pair['People'], background_image_path)
    
    # Clean up temporary background image file if created
    if background_image:
        cleanup_temp_file(background_image_path)
    
    # Save the presentation
    save_presentation(prs, output_file)
    return f"Created {len(pairs_list)} Call to Worship slides"


def add_call_to_worship_slide(prs, leader_text, people_text, background_image_path=None):
    """Add a Call to Worship slide with Leader/People format"""
    # Sanitize input to remove any accidental "Leader:" or "People:" prefixes
    if leader_text.strip().lower().startswith("leader:"):
        leader_text = leader_text.split(':', 1)[-1].lstrip()
    if people_text.strip().lower().startswith("people:"):
        people_text = people_text.split(':', 1)[-1].lstrip()
        
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    set_slide_background(slide, background_image_path)
    
    # Add flower placeholder image in top right corner
    placeholder_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))),
        'public', 'images', 'placeholders', 'flower-placholder.jpg'
    )
    
    if os.path.exists(placeholder_path):
        # Position: 9.33 inches horizontal, 0.25 inches vertical
        left = Inches(9.33)
        top = Inches(0.25)
        # Size: width 3.72 inches, height 2 inches
        width = Inches(3.72)
        height = Inches(2)
        slide.shapes.add_picture(placeholder_path, left, top, width=width, height=height)
    
    # Add title - LEFT ALIGNED
    title_left = Inches(0.26)  # Match hymn slide positioning
    title_top = Inches(0.03)  # Position at 0.03 inches from top
    title_width = Inches(8.91)  # Match hymn slide width
    title_height = Inches(1.2)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "CALL TO WORSHIP"
    title_frame.word_wrap = True
    title_frame.margin_left = Inches(0)
    title_frame.margin_right = Inches(0)
    
    # Format title - LEFT ALIGNED, BOLD and UNDERLINED
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "Arial Narrow"
    title_para.font.size = Pt(54)  # Set to 54pt
    title_para.font.bold = True
    title_para.font.underline = True  # Add underline
    title_para.font.color.rgb = RGBColor(0, 0, 0)  # Black text
    title_para.alignment = PP_ALIGN.LEFT  # Changed to LEFT
    
    # Add white glow effect to title
    add_text_glow(title_para, glow_radius=6, color_rgb=(255, 255, 255))
    
    # Single text box for both Leader and People sections
    font_size = 50
    content_left = Inches(0)
    content_width = Inches(13.33)
    content_top = Inches(2.5)  # Moved higher to 2.5 inches from top
    content_height = Inches(5)  # Increased height for both sections
    
    # Create single text box (no background highlight)
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.TOP
    content_frame.margin_left = Inches(0.5)
    content_frame.margin_right = Inches(0.5)
    content_frame.clear()
    
    # Leader paragraph (split into runs to better control glow like template)
    leader_para = content_frame.paragraphs[0]
    leader_para.text = ""
    leader_para.font.name = "Arial Narrow"
    leader_para.font.size = Pt(font_size)
    leader_para.font.color.rgb = RGBColor(0, 0, 0)
    leader_para.alignment = PP_ALIGN.LEFT
    leader_para.font.bold = True
    leader_para.level = 0
    leader_para.space_after = Pt(12)  # Add space after Leader section
    leader_para.space_before = Pt(0)
    leader_para.left_indent = Inches(2.2)
    leader_para.first_line_indent = Inches(-2.2)
    set_hanging_indent(leader_para, left_inches=2.2)
    # Build runs: label gets white glow, text gets scheme glow
    leader_label = leader_para.add_run()
    leader_label.text = "Leader: "
    add_run_glow(leader_label, color_rgb=(255, 255, 255), glow_radius_pt=6)
    leader_text_run = leader_para.add_run()
    leader_text_run.text = leader_text
    # Use white glow for Leader text (medium radius) per style guide
    add_run_glow(leader_text_run, color_rgb=(255, 255, 255), glow_radius_pt=10)
    
    # Add People paragraph (using add_paragraph for the second paragraph)
    people_para = content_frame.add_paragraph()
    people_para.text = ""
    people_para.font.name = "Arial Narrow"
    people_para.font.size = Pt(font_size)
    people_para.font.color.rgb = RGBColor(0, 0, 0)
    people_para.alignment = PP_ALIGN.LEFT
    people_para.font.bold = True
    people_para.level = 0
    people_para.space_after = Pt(0)
    people_para.space_before = Pt(12)  # Add space before People section
    people_para.left_indent = Inches(2.2)
    people_para.first_line_indent = Inches(-2.2)
    set_hanging_indent(people_para, left_inches=2.2)
    # Build runs like template
    people_label = people_para.add_run()
    people_label.text = "People: "
    # Apply fill, effects, then highlight (order matters to avoid glow around highlight)
    add_run_fill(people_label, color_rgb=(0, 0, 0))
    # Combine white glow and outline for text visibility
    set_run_effects(
        people_label,
        glow_rgb=(255, 255, 255),
        glow_radius_pt=6  # Smaller glow for label
    )
    add_run_outline(people_label, color_rgb=(255, 255, 255), width_pt=1.0)
    set_run_highlight(people_label, color_rgb=(255, 255, 0))  # Yellow background last
    
    people_text_run = people_para.add_run()
    people_text_run.text = people_text
    # Apply fill, effects, then highlight
    add_run_fill(people_text_run, color_rgb=(0, 0, 0))
    # Use white glow for better text visibility
    set_run_effects(
        people_text_run,
        glow_rgb=(255, 255, 255),
        glow_radius_pt=10  # Larger glow for main text
    )
    add_run_outline(people_text_run, color_rgb=(255, 255, 255), width_pt=1.0)
    set_run_highlight(people_text_run, color_rgb=(255, 255, 0))  # Yellow background last
    # Also set endParaRPr effects similar to template for better renderer support
    add_end_paragraph_glow_and_highlight(people_para, scheme="bg1", glow_radius_pt=10.0, highlight_rgb=(255, 255, 0))


@router.post("/generate-call-to-worship")
async def generate_call_to_worship_endpoint(request: CallToWorshipRequest):
    """Generate Call to Worship PowerPoint slides"""
    try:
        # Extract pairs and background info
        pairs = request.pairs if request.pairs else []
        background_image = request.background_image
        
        if not pairs:
            # If no pairs provided, try to parse from text
            text = request.text if request.text else ''
            if not text:
                raise HTTPException(status_code=400, detail='No pairs or text provided')
            
            # Parse text into Leader/People format
            lines = text.split('\n')
            leader_text = ""
            people_text = ""
            
            for line in lines:
                line = line.strip()
                if line.lower().startswith('leader:'):
                    leader_text = line[7:].strip()
                elif line.lower().startswith('people:'):
                    people_text = line[7:].strip()
            
            if leader_text or people_text:
                pairs = [{'Leader': leader_text, 'People': people_text}]
            else:
                # Use the whole text as a single slide
                pairs = [{'Leader': text, 'People': ''}]
        
        # Create temporary file for the PowerPoint
        output_path = create_temp_file(suffix='.pptx')
        
        # Use background image if provided (base64 only)
        background = background_image
        
        # Generate the PowerPoint
        result = create_call_to_worship_slides_from_dict(pairs, output_path, background)
        
        # Return the file
        return FileResponse(
            path=output_path,
            filename="call_to_worship.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))