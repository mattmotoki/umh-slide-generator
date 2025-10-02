import os
import base64
import tempfile
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.util import Inches
from pptx.dml.color import RGBColor


# Base presentation functions

def create_presentation(width=Inches(13.33), height=Inches(7.5)):
    """Create a new PowerPoint presentation with specified dimensions"""
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen by default)
    prs.slide_width = width
    prs.slide_height = height
    
    return prs


def save_presentation(prs, output_file):
    """Save presentation to file"""
    prs.save(output_file)
    return output_file


def process_background_image(background_image_data):
    """
    Process base64 background image data and save to temporary file.
    Returns path to temporary image file.
    """
    if not background_image_data:
        return None
    
    try:
        # Remove data URL prefix if present (data:image/jpeg;base64,...)
        if ',' in background_image_data:
            background_image_data = background_image_data.split(',')[1]
        
        # Decode base64 data
        image_data = base64.b64decode(background_image_data)
        
        # Create temporary file for the image
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            tmp_file.write(image_data)
            return tmp_file.name
            
    except Exception as e:
        print(f"Error processing background image: {e}")
        return None


# Removed deprecated get_default_background_path function
# Background images should always be sent from the frontend as base64


def set_slide_background(slide, background_image_path=None):
    """Set slide background to either an image or white"""
    if background_image_path and os.path.exists(background_image_path):
        try:
            # Add background image
            left = 0
            top = 0
            # Access presentation through slide's parent
            prs = slide.part.package.presentation_part.presentation
            slide.shapes.add_picture(background_image_path, left, top, 
                                   prs.slide_width, 
                                   prs.slide_height)
            return True
        except Exception as e:
            print(f"Error adding background image: {e}")
    
    # Fall back to white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    return False


def cleanup_temp_file(file_path):
    """Clean up temporary file if it exists and is in temp directory"""
    if file_path and os.path.exists(file_path):
        # Only delete if it's a temporary file (contains 'tmp' in the path)
        if 'tmp' in file_path or file_path.startswith(tempfile.gettempdir()):
            try:
                os.unlink(file_path)
            except Exception as e:
                print(f"Error cleaning up temporary file: {e}")


# Text effect functions

def add_text_glow(paragraph, glow_radius: int = 6, color_rgb: tuple[int, int, int] = (255, 255, 255)) -> None:
    """Add a white (or provided color) glow effect to all runs in a paragraph.

    The effect is applied via XML modifications because python-pptx does not
    currently expose glow styling via its public API.

    Args:
        paragraph: The paragraph whose runs should receive a glow effect.
        glow_radius: Glow radius in points.
        color_rgb: (R, G, B) color tuple for the glow.
    """
    for run in paragraph.runs:
        run_properties = run._r.get_or_add_rPr()
        glow_effect = parse_xml(
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:glow rad="{int(glow_radius * 12700)}">'
            f'    <a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}">'
            f'      <a:alpha val="100000"/>'
            f'    </a:srgbClr>'
            f'  </a:glow>'
            f'</a:effectLst>'
        )
        run_properties.append(glow_effect)



def add_text_highlight(paragraph, color_rgb: tuple[int, int, int] = (255, 255, 0)) -> None:
    """Apply text highlight color to all runs in a paragraph via XML.

    PowerPoint supports run highlighting using the <a:highlight> element under
    run properties (a:rPr). python-pptx does not expose this directly.
    """
    for run in paragraph.runs:
        run_properties = run._r.get_or_add_rPr()
        highlight_xml = parse_xml(
            f'<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"/>'
            f'</a:highlight>'
        )
        run_properties.append(highlight_xml)


def set_hanging_indent(paragraph, left_inches: float = 2.2) -> None:
    """Set a hanging indent on the paragraph: first line at 0, wraps at left_inches.

    Implemented by setting paragraph properties marL (left margin) to left_inches
    and indent to -left_inches in EMUs.
    """
    try:
        pPr = paragraph._p.get_or_add_pPr()
        left_emu = int(Inches(left_inches))
        pPr.set('marL', str(left_emu))
        pPr.set('indent', str(-left_emu))
    except Exception:
        # Best-effort; ignore if underlying XML structure differs
        pass

def add_text_outline(paragraph, color_rgb: tuple[int, int, int] = (255, 255, 255), width_pt: float = 1.0) -> None:
    """Add a text outline to all runs in a paragraph via XML.

    Outline helps the text stand out on colored backgrounds.
    """
    width = max(1, int(width_pt * 12700))
    for run in paragraph.runs:
        run_properties = run._r.get_or_add_rPr()
        outline_xml = parse_xml(
            f'<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width}">'
            f'  <a:solidFill><a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"/></a:solidFill>'
            f'</a:ln>'
        )
        run_properties.append(outline_xml)


def add_text_glow_scheme(paragraph, scheme: str = "bg1", glow_radius_pt: float = 10.0) -> None:
    """Add glow using a scheme color (e.g., bg1) for all runs in a paragraph."""
    rad = max(1, int(glow_radius_pt * 12700))
    for run in paragraph.runs:
        run_properties = run._r.get_or_add_rPr()
        glow_effect = parse_xml(
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:glow rad="{rad}">'
            f'    <a:schemeClr val="{scheme}"/>'
            f'  </a:glow>'
            f'</a:effectLst>'
        )
        run_properties.append(glow_effect)


def add_end_paragraph_glow_and_highlight(paragraph, scheme: str = "bg1", glow_radius_pt: float = 10.0,
                                         highlight_rgb: tuple[int, int, int] | None = None) -> None:
    """Attach glow (scheme color) and optional highlight to endParaRPr for a paragraph."""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        endRPr = None
        # Try to find existing endParaRPr
        for child in pPr:
            if child.tag.endswith('endParaRPr'):
                endRPr = child
                break
        if endRPr is None:
            endRPr = parse_xml('<a:endParaRPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            pPr.append(endRPr)
        # Add glow
        rad = max(1, int(glow_radius_pt * 12700))
        effectLst = parse_xml(
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:glow rad="{rad}"><a:schemeClr val="{scheme}"/></a:glow>'
            f'</a:effectLst>'
        )
        endRPr.append(effectLst)
        # Optional highlight
        if highlight_rgb is not None:
            highlight = parse_xml(
                f'<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'  <a:srgbClr val="{highlight_rgb[0]:02X}{highlight_rgb[1]:02X}{highlight_rgb[2]:02X}"/>'
                f'</a:highlight>'
            )
            endRPr.append(highlight)
    except Exception:
        pass


def add_run_glow(run, color_rgb: tuple[int, int, int] | None = (255, 255, 255), scheme: str | None = None,
                 glow_radius_pt: float = 10.0) -> None:
    """Add glow to a single run using either srgb color or scheme color."""
    rad = max(1, int(glow_radius_pt * 12700))
    run_properties = run._r.get_or_add_rPr()
    if scheme:
        glow_effect = parse_xml(
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:glow rad="{rad}"><a:schemeClr val="{scheme}"/></a:glow>'
            f'</a:effectLst>'
        )
    else:
        glow_effect = parse_xml(
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:glow rad="{rad}"><a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"><a:alpha val="100000"/></a:srgbClr></a:glow>'
            f'</a:effectLst>'
        )
    run_properties.append(glow_effect)


def add_run_highlight(run, color_rgb: tuple[int, int, int] = (255, 255, 0)) -> None:
    """Add highlight to a single run via XML."""
    run_properties = run._r.get_or_add_rPr()
    highlight_xml = parse_xml(
        f'<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"/>'
        f'</a:highlight>'
    )
    run_properties.append(highlight_xml)


def add_run_outline(run, color_rgb: tuple[int, int, int] = (255, 255, 255), width_pt: float = 0.75) -> None:
    """Add a text outline to a single run via XML."""
    width = max(1, int(width_pt * 12700))
    run_properties = run._r.get_or_add_rPr()
    outline_xml = parse_xml(
        f'<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width}">'
        f'  <a:solidFill><a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"/></a:solidFill>'
        f'</a:ln>'
    )
    run_properties.append(outline_xml)


def add_run_fill(run, color_rgb: tuple[int, int, int] = (0, 0, 0)) -> None:
    """Set solid text fill color on a single run via XML."""
    run_properties = run._r.get_or_add_rPr()
    fill_xml = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"/>'
        f'</a:solidFill>'
    )
    run_properties.append(fill_xml)


def add_run_outer_shadow(
    run,
    color_rgb: tuple[int, int, int] = (255, 255, 255),
    blur_radius_pt: float = 2.0,
    distance_pt: float = 0.5,
    direction_degrees: int = 0,
) -> None:
    """Add a subtle outer shadow to a single run via XML (can simulate an outline/edge).

    Note: PowerPoint uses EMUs for sizes; 1pt ≈ 12700 EMUs. Direction is in degrees * 60000 (OOXML units).
    """
    blur = max(1, int(blur_radius_pt * 12700))
    dist = max(0, int(distance_pt * 12700))
    dir_ooxml = int(direction_degrees * 60000)
    run_properties = run._r.get_or_add_rPr()
    shadow_xml = parse_xml(
        f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <a:outerShdw blurRad="{blur}" dist="{dist}" dir="{dir_ooxml}" algn="ctr" rotWithShape="0">'
        f'    <a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"><a:alpha val="100000"/></a:srgbClr>'
        f'  </a:outerShdw>'
        f'</a:effectLst>'
    )
    run_properties.append(shadow_xml)


def set_run_effects(
    run,
    glow_rgb: tuple[int, int, int] | None = None,
    glow_radius_pt: float = 10.0,
    inner_shadow_rgb: tuple[int, int, int] | None = None,
    inner_blur_radius_pt: float = 1.5,
    inner_distance_pt: float = 0.0,
    outer_shadow_rgb: tuple[int, int, int] | None = None,
    blur_radius_pt: float = 1.5,
    distance_pt: float = 0.25,
    direction_degrees: int = 0,
) -> None:
    """Replace any existing effectLst with a single combined list (glow + outer shadow).

    Order: glow first, then outerShdw inside the same a:effectLst.
    """
    rPr = run._r.get_or_add_rPr()
    # Remove existing effectLst children
    to_remove = []
    for child in list(rPr):
        if child.tag.endswith('effectLst'):
            to_remove.append(child)
    for node in to_remove:
        rPr.remove(node)

    parts: list[str] = []
    if glow_rgb is not None:
        rad = max(1, int(glow_radius_pt * 12700))
        parts.append(
            f'<a:glow rad="{rad}"><a:srgbClr val="{glow_rgb[0]:02X}{glow_rgb[1]:02X}{glow_rgb[2]:02X}"><a:alpha val="100000"/></a:srgbClr></a:glow>'
        )
    if inner_shadow_rgb is not None:
        iblur = max(1, int(inner_blur_radius_pt * 12700))
        idist = max(0, int(inner_distance_pt * 12700))
        idir = int(direction_degrees * 60000)
        parts.append(
            f'<a:innerShdw blurRad="{iblur}" dist="{idist}" dir="{idir}">' 
            f'  <a:srgbClr val="{inner_shadow_rgb[0]:02X}{inner_shadow_rgb[1]:02X}{inner_shadow_rgb[2]:02X}"><a:alpha val="100000"/></a:srgbClr>'
            f'</a:innerShdw>'
        )
    if outer_shadow_rgb is not None:
        blur = max(1, int(blur_radius_pt * 12700))
        dist = max(0, int(distance_pt * 12700))
        dir_ooxml = int(direction_degrees * 60000)
        parts.append(
            f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{dir_ooxml}" algn="ctr" rotWithShape="0">'
            f'  <a:srgbClr val="{outer_shadow_rgb[0]:02X}{outer_shadow_rgb[1]:02X}{outer_shadow_rgb[2]:02X}"><a:alpha val="100000"/></a:srgbClr>'
            f'</a:outerShdw>'
        )
    if parts:
        effect_xml = parse_xml(
            '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">' + ''.join(parts) + '</a:effectLst>'
        )
        rPr.append(effect_xml)


def set_run_highlight(run, color_rgb: tuple[int, int, int] = (255, 255, 0)) -> None:
    """Replace any existing highlight on a run and append a new one (after effects)."""
    rPr = run._r.get_or_add_rPr()
    # Remove existing highlight children
    to_remove = []
    for child in list(rPr):
        if child.tag.endswith('highlight'):
            to_remove.append(child)
    for node in to_remove:
        rPr.remove(node)
    # Append new highlight
    highlight_xml = parse_xml(
        f'<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <a:srgbClr val="{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}"/>'
        f'</a:highlight>'
    )
    rPr.append(highlight_xml)


def set_textframe_no_autofit(text_frame, left_inches: float | None = None, right_inches: float | None = None) -> None:
    """Set a text frame to noAutofit and optionally adjust left/right insets."""
    try:
        tx = text_frame._txBody
        bodyPr = tx.get_or_add_bodyPr()
        # Remove spAutoFit if present
        for child in list(bodyPr):
            if child.tag.endswith('spAutoFit'):
                bodyPr.remove(child)
            if child.tag.endswith('noAutofit'):
                # already present
                pass
        # Add noAutofit
        no_autofit = parse_xml('<a:noAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        bodyPr.append(no_autofit)
        # Set insets
        if left_inches is not None:
            bodyPr.set('lIns', str(int(Inches(left_inches))))
        if right_inches is not None:
            bodyPr.set('rIns', str(int(Inches(right_inches))))
    except Exception:
        pass

